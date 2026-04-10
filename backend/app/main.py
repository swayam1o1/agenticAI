import os
import json as _json

from fastapi import FastAPI, UploadFile, File, Form, HTTPException, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from typing import List, Optional
import uvicorn
import io
import PyPDF2
import pptx

from .memory import FAISSMemory
from .agent import StudyAgent
from .models import AgentRequest, AgentResponse, QuizAnswerSubmission, TaskStatusUpdate
from .storage import Storage
from .orchestrator import AgenticOrchestrator
from .learn_orchestrator import LearnOrchestrator

ROOT_DIR = os.path.dirname(os.path.dirname(__file__))
DATA_DIR = os.path.join(ROOT_DIR, "data")
os.makedirs(DATA_DIR, exist_ok=True)
CHAT_DB = os.path.join(DATA_DIR, "chat.db")
OLLAMA_MODEL = "llama3.2"
EMBED_MODEL = "mxbai-embed-large"

app = FastAPI(title="Agentic Study Buddy", version="0.1.0")

# CORS (adjust origins as needed)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)

memory = FAISSMemory(data_dir=DATA_DIR, embed_model=EMBED_MODEL)
storage = Storage(f"sqlite:///{CHAT_DB}")
agent = StudyAgent(memory=memory, model=OLLAMA_MODEL, storage=storage)
orchestrator = AgenticOrchestrator(agent=agent, storage=storage, memory=memory)
learn_orchestrator = LearnOrchestrator(agent=agent, storage=storage, memory=memory)


@app.get("/api/health")
def health():
    return {"status": "ok"}


def _memory_analysis_loop(session_id: str):
    try:
        agent.run(task="mindmap", user_input="", history=[], session_id=session_id)
    except Exception as e:
        print(f"Memory loop failed: {e}")

@app.post("/api/memory")
async def ingest_memory(
    background_tasks: BackgroundTasks, 
    texts: Optional[List[str]] = Form(default=None), 
    file: Optional[UploadFile] = File(default=None),
    session_id: Optional[str] = Form(default=None)
):
    payload_texts: List[str] = []
    if texts:
        payload_texts.extend(texts)
    if file is not None:
        filename = file.filename.lower() if hasattr(file, 'filename') and file.filename else ""
        content_bytes = await file.read()
        
        if filename.endswith(".pdf"):
            try:
                reader = PyPDF2.PdfReader(io.BytesIO(content_bytes))
                text = ""
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                payload_texts.append(text)
            except Exception as e:
                payload_texts.append(f"Failed to read PDF: {e}")
        elif filename.endswith(".pptx"):
            try:
                prs = pptx.Presentation(io.BytesIO(content_bytes))
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                payload_texts.append(text)
            except Exception as e:
                payload_texts.append(f"Failed to read PPTX: {e}")
        else:
            content = content_bytes.decode("utf-8", errors="ignore")
            payload_texts.append(content)
            
    if not payload_texts:
        return {"added": 0, "ids": []}
    ids = memory.add_texts(payload_texts)
    
    if session_id:
        background_tasks.add_task(_memory_analysis_loop, session_id)
        
    return {"added": len(ids), "ids": ids}


@app.get("/api/memory")
async def read_memory_bank(limit: int = 100):
    return {"items": memory.list_memories(limit=limit)}


@app.post("/api/agent", response_model=AgentResponse)
async def run_agent(req: AgentRequest):
    history = [m.dict() for m in (req.history or [])]
    result = agent.run(task=req.task, user_input=req.input, history=history, session_id=req.session_id)
    session_id = result.get("session_id")
    meta = dict(result.get("meta", {}))
    if session_id:
        meta["next_action"] = orchestrator.get_next_recommended_action(session_id)
    response = AgentResponse(task=req.task, output=result["output"], meta=meta)
    response.session_id = session_id
    return response

def _extract_requested_topics(history: list) -> list:
    """Extract topics the user explicitly asked to learn about, bypassing LLM."""
    import re
    patterns = [
        r"(?:please\s+)?teach\s+me\s+about\s+(.+?)(?:\.\s*$|$)",
        r"(?:please\s+)?explain\s+(?:to\s+me\s+)?(.+?)(?:\.\s*$|$)",
        r"what\s+is\s+(?:a\s+|an\s+|the\s+)?(.+?)(?:\?\s*$|$)",
        r"how\s+does\s+(.+?)\s+work",
        r"tell\s+me\s+about\s+(.+?)(?:\.\s*$|$)",
        r"i\s+want\s+to\s+(?:learn|understand)\s+(.+?)(?:\.\s*$|$)",
    ]
    seen = set()
    topics = []
    for msg in history:
        if msg.get("role") != "user":
            continue
        content = msg.get("content", "").strip()
        for pattern in patterns:
            m = re.search(pattern, content, re.IGNORECASE)
            if m:
                topic = m.group(1).strip().rstrip("?.,!").strip()
                # Filter out generic filler
                if 2 < len(topic) < 80 and topic.lower() not in seen:
                    seen.add(topic.lower())
                    topics.append(topic)
                break
    return topics


def background_analysis_process(session_id: str):
    """Run analysis and auto-quiz generation in the background"""
    try:
        history = storage.get_history(session_id)
        user_msgs = [m for m in history if m.get("role") == "user"]
        if len(user_msgs) == 0:
            return

        # --- STEP 1: Directly extract and write topics WITHOUT the LLM ---
        requested_topics = _extract_requested_topics(history)
        if requested_topics:
            # Format as a numbered list that log_weak_topics can parse
            topic_summary = "\n".join(
                f"{i+1}. {t}: The learner explicitly requested to learn this topic."
                for i, t in enumerate(requested_topics[:5])
            )
            print(f"[Background] Direct-injecting {len(requested_topics)} topics: {requested_topics}")
            storage.log_weak_topics(session_id, topic_summary)

        # --- STEP 2: Run LLM-based analysis for deeper pattern detection ---
        print(f"[Background] Real-time chat analysis for session {session_id} ({len(user_msgs)} user messages)")
        # silent=True prevents "chat analysis" being logged as a fake user message
        agent.run(task="analyze", user_input="chat analysis", history=history, session_id=session_id, silent=True)
        
        # Recalibrate the high-level roadmap document alongside local tasks
        print(f"[Background] Syncing roadmap plan for session {session_id}")
        agent.run(task="roadmap", user_input="Recalibrate based on recent chat", history=history, session_id=session_id, silent=True)

        # Check if quiz is recommended and auto-generate
        next_action = orchestrator.get_next_recommended_action(session_id)
        if next_action and next_action.get("action") == "quiz":
            print(f"[Background] Auto-generating quiz for session {session_id}")
            agent.run(task="quiz", user_input="Auto-generated focused quiz", history=[], session_id=session_id, silent=True)
    except Exception as e:
        print(f"[Background Error] {e}")


@app.post("/api/tutor/stream")
async def stream_tutor(req: AgentRequest, background_tasks: BackgroundTasks):
    history = [m.dict() for m in (req.history or [])]
    session_id = storage.ensure_session(req.session_id)
    storage.log_message(session_id, "user", req.input, task="tutor")

    # Get context from memory
    hits = memory.similarity_search(req.input, k=5)
    ctx = "\n\n".join([f"[Score {score:.2f}] {md['text']}" for _, score, md in hits])

    # Build conversation history text
    history_text = ""
    if history:
        recent = history[-10:]
        history_text = "\n".join([f"{m.get('role','').upper()}: {m.get('content','')}" for m in recent])

    proactive = orchestrator.orchestrate_learning_cycle(session_id, "tutor")
    weak_topics = proactive.get("weak_topics", [])
    roadmap_focus = proactive.get("roadmap_tasks", [])

    prompt = (
        "You are an expert AI tutor. Your teaching style:\n"
        "- Explain concepts step-by-step with clear structure\n"
        "- Use real-world analogies and examples\n"
        "- Format with markdown: ## headers, **bold** key terms, bullet points\n"
        "- Use code blocks for code or pseudocode\n"
        "- Be concise but thorough, ask follow-up questions when helpful\n\n"
    )
    if ctx.strip():
        prompt += f"Reference Material:\n{ctx}\n\n"
    if history_text:
        prompt += f"Conversation so far:\n{history_text}\n\n"
    if weak_topics:
        prompt += "Known weak areas to reinforce:\n"
        prompt += "\n".join([f"- {item.get('title', 'Unknown')}: {item.get('detail', '')}" for item in weak_topics])
        prompt += "\n\n"
    if roadmap_focus:
        prompt += "Open study tasks:\n"
        prompt += "\n".join([f"- {task.get('title', 'Task')}: {task.get('detail', '')}" for task in roadmap_focus])
        prompt += "\n\n"
    prompt += f"Student: {req.input}\nTutor:"

    full_tokens: List[str] = []

    async def event_stream():
        async for token in agent.llm.generate_stream_async(prompt):
            full_tokens.append(token)
            yield f"data: {_json.dumps({'token': token})}\n\n"
        answer = "".join(full_tokens)
        storage.log_message(session_id, "assistant", answer, task="tutor")
        memory.add_texts(
            [
                (
                    f"Task: tutor\n"
                    f"Learner request: {req.input.strip()}\n"
                    f"Tutor response summary: {answer[:1200]}\n"
                    f"Referenced memory chunks: {len(hits)}"
                )
            ],
            [{"session_id": session_id, "task": "tutor", "kind": "interaction-summary"}],
        )
        next_action = orchestrator.get_next_recommended_action(session_id)
        yield f"data: {_json.dumps({'done': True, 'session_id': session_id, 'next_action': next_action})}\n\n"

    background_tasks.add_task(background_analysis_process, session_id)
    return StreamingResponse(
        event_stream(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


@app.get("/api/history")
async def read_history(session_id: str):
    if not session_id:
        raise HTTPException(status_code=400, detail="session_id is required")
    return {"session_id": session_id, "messages": storage.get_history(session_id)}


@app.get("/api/weak-topics")
async def read_weak_topics(session_id: str):
    if not session_id:
        raise HTTPException(status_code=400, detail="session_id is required")
    return {"session_id": session_id, "weak_topics": storage.get_weak_topics(session_id)}


@app.get("/api/analysis")
async def get_analysis(session_id: str):
    """Get the latest analysis summary for a session"""
    if not session_id:
        raise HTTPException(status_code=400, detail="session_id is required")
    # Get the most recent analysis message
    history = storage.get_history(session_id)
    for msg in reversed(history):
        if msg.get("task") == "analyze" and msg.get("role") == "assistant":
            return {"session_id": session_id, "summary": msg.get("content", ""), "timestamp": msg.get("timestamp")}
    return {"session_id": session_id, "summary": None}


@app.get("/api/quiz-history")
async def read_quiz_history(session_id: str):
    if not session_id:
        raise HTTPException(status_code=400, detail="session_id is required")
    return {"session_id": session_id, "quiz_history": storage.get_quiz_history(session_id)}


@app.get("/api/roadmap")
async def read_roadmap(session_id: str):
    if not session_id:
        raise HTTPException(status_code=400, detail="session_id is required")
    return {"session_id": session_id, "tasks": storage.get_roadmap_tasks(session_id)}


@app.post("/api/roadmap/task-status")
async def update_task_status(payload: TaskStatusUpdate):
    if not storage.update_task_status(payload.session_id, payload.task_id, payload.status):
        raise HTTPException(status_code=400, detail="Task not found or invalid status")
    return {"status": "ok"}


def _quiz_analysis_loop(session_id: str):
    try:
        agent.run(task="analyze", user_input="quiz", history=[], session_id=session_id)
        agent.run(task="roadmap", user_input="", history=[], session_id=session_id)
    except Exception as e:
        print(f"Quiz auto loop failed: {e}")

@app.post("/api/quiz-answer")
async def submit_quiz_answer(payload: QuizAnswerSubmission, background_tasks: BackgroundTasks):
    if not storage.record_quiz_answer(
        payload.session_id,
        payload.attempt_id,
        payload.question_id,
        payload.selected_index,
        payload.selected_option,
        payload.is_correct,
        note=payload.note,
        confidence=payload.confidence,
    ):
        raise HTTPException(status_code=400, detail="Failed to record answer")
    background_tasks.add_task(_quiz_analysis_loop, payload.session_id)
    return {"status": "ok"}


@app.get("/api/recommendations")
async def get_recommendations(session_id: str):
    """Get AI-powered learning recommendations based on progress"""
    if not session_id:
        raise HTTPException(status_code=400, detail="session_id is required")
    try:
        next_action = orchestrator.get_next_recommended_action(session_id)
        performance = orchestrator.analyze_quiz_performance(session_id)
        return {
            "session_id": session_id,
            "next_action": next_action,
            "performance": performance
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/learn/start")
async def start_learning_concept(payload: dict):
    """Start learning a concept - Phase 1: Teaching"""
    session_id = payload.get("session_id")
    concept = payload.get("concept")
    if not session_id or not concept:
        raise HTTPException(status_code=400, detail="session_id and concept are required")
    try:
        result = learn_orchestrator.start_learning(session_id, concept)
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/learn/quiz")
async def generate_learning_quiz(payload: dict):
    """Generate quiz for concept - Phase 2: Quiz"""
    session_id = payload.get("session_id")
    concept = payload.get("concept")
    focus_weak = payload.get("focus_weak_areas", False)
    if not session_id or not concept:
        raise HTTPException(status_code=400, detail="session_id and concept are required")
    try:
        result = learn_orchestrator.generate_concept_quiz(session_id, concept, focus_weak)
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/learn/analyze")
async def analyze_learning_quiz(payload: dict):
    """Analyze quiz results - Phase 3: Analysis"""
    session_id = payload.get("session_id")
    attempt_id = payload.get("attempt_id")
    concept = payload.get("concept")
    if not session_id or not attempt_id or not concept:
        raise HTTPException(status_code=400, detail="session_id, attempt_id, and concept are required")
    try:
        result = learn_orchestrator.analyze_quiz_results(session_id, attempt_id, concept)
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/learn/progress")
async def get_learning_progress(session_id: str, concept: Optional[str] = None):
    """Get learning progress for concept(s)"""
    if not session_id:
        raise HTTPException(status_code=400, detail="session_id is required")
    try:
        result = learn_orchestrator.get_learning_progress(session_id, concept)
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/mastery")
async def get_concept_mastery(session_id: str):
    """Get all concept mastery data"""
    if not session_id:
        raise HTTPException(status_code=400, detail="session_id is required")
    return {"session_id": session_id, "masteries": storage.get_concept_mastery(session_id)}


@app.post("/api/synthesize")
async def generate_study_guide(payload: dict):
    """Generate a NotebookLM-style comprehensive Study Guide"""
    session_id = payload.get("session_id")
    if not session_id:
        raise HTTPException(status_code=400, detail="session_id is required")
    try:
        result = agent.run(task="synthesize", user_input="", history=[], session_id=session_id)
        guide = result.get("output", {}).get("guide", "Failed to generate guide")
        return {"session_id": session_id, "guide": guide}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/mindmap")
async def generate_mindmap(session_id: str, topic: Optional[str] = None):
    """Generate a Mermaid syntax mind map of the memory bank / weak topics"""
    if not session_id:
        raise HTTPException(status_code=400, detail="session_id is required")
    try:
        user_in = topic if topic else ""
        result = agent.run(task="mindmap", user_input=user_in, history=[], session_id=session_id)
        mindmap_raw = result.get("output", {}).get("mindmap", "mindmap\n  root((Failed to generate map))")
        
        # 🛡️ Algorithmic Mermaid Compiler
        # Converts arbitrary LLM syntax into strict id("text") shapes to guarantee 100% render safety
        import re
        clean_lines = ["mindmap"]
        node_idx = 0

        for line in mindmap_raw.splitlines():
            line = line.replace("\t", "  ")
            stripped = line.strip()
            
            if not stripped or stripped.lower() == "mindmap":
                continue
                
            indent = len(line) - len(stripped)
            
            # A valid Mermaid tree node MUST be indented. 0-indent paragraphs are just LLM hallucinated chat prefixes.
            if indent == 0:
                continue
                
            node_idx += 1
            
            # 🛡️ Ultra-Stable Mermaid Compiler
            # Standard mindmap syntax:  "Label"
            # Scrub all characters that could be misinterpreted as shape markers or syntax breaks
            safe_text = stripped.replace('"', "'")
            for char in "()[]{}:;":
                safe_text = safe_text.replace(char, "")
                
            # If the LLM hallucinated its own Node IDs or shape markers at the start, strip them
            safe_text = re.sub(r'^id\d+\s*', '', safe_text).strip()
            safe_text = re.sub(r'^[\[\{\(]+|[\]\}\)]+$', '', safe_text).strip()
            
            # Use simple indentation + quoted text. This is the most resilient Mermaid Mindmap format.
            clean_lines.append(f"{' ' * indent}\"{safe_text}\"")
            
        mindmap = "\n".join(clean_lines)

        if mindmap and "id1" in mindmap:
             storage.log_mindmap(session_id, topic, mindmap)
        return {"session_id": session_id, "mindmap": mindmap}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/mindmap-history")
async def get_mindmap_history(session_id: str):
    if not session_id:
        raise HTTPException(status_code=400, detail="session_id is required")
    try:
        history = storage.get_mindmap_history(session_id)
        return {"session_id": session_id, "history": history}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/flashcards")
async def generate_flashcards(session_id: str, topic: Optional[str] = None):
    """Generate dynamic flashcards based on weak areas"""
    if not session_id:
        raise HTTPException(status_code=400, detail="session_id is required")
    try:
        user_in = topic if topic else ""
        result = agent.run(task="flashcards", user_input=user_in, history=[], session_id=session_id)
        cards = result.get("output", {}).get("flashcards", [])
        if cards:
            storage.log_flashcards(session_id, topic, cards)
        return {"session_id": session_id, "flashcards": cards}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/flashcards-history")
async def get_flashcards_history(session_id: str):
    if not session_id:
        raise HTTPException(status_code=400, detail="session_id is required")
    try:
        history = storage.get_flashcards_history(session_id)
        return {"session_id": session_id, "history": history}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


if __name__ == "__main__":
    uvicorn.run("app.main:app", host="127.0.0.1", port=8001, reload=True)
